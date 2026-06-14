"""Attachment text extraction for the Live Inbox.

Fail-and-fix-early: never silently drop an attachment. Extract its text (txt/md/csv/
pdf/docx/...) so the agent + RAG see it, or — if we can't read it — return None so the
caller NOTES it explicitly in the message ("could not extract …") rather than hiding it.
"""
import io
import logging

log = logging.getLogger("attachments")

TEXT_EXT = (".txt", ".md", ".csv", ".tsv", ".vtt", ".eml", ".json", ".log", ".yaml", ".yml")


def _txt(data):
    return data.decode("utf-8", "ignore")


def _pdf(data):
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(data))
    return "\n".join((page.extract_text() or "") for page in reader.pages)


def _docx(data):
    import docx
    document = docx.Document(io.BytesIO(data))
    return "\n".join(p.text for p in document.paragraphs)


def _pptx(data):
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    out = []
    for i, slide in enumerate(prs.slides, 1):
        out.append("--- slide %d ---" % i)
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                out.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    out.append(" | ".join(c.text for c in row.cells))
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes and notes.strip():
                out.append("[notes] " + notes)
    return "\n".join(out)


def extract(filename, content_type, data):
    """Text from one attachment, or None if unsupported/unreadable (caller must note it)."""
    fn = (filename or "").lower()
    ct = (content_type or "").lower()
    try:
        import transcribe
        if transcribe.is_media(fn, ct):
            return transcribe.transcribe(filename, data, content_type)
        if ct.startswith("text/") or fn.endswith(TEXT_EXT):
            return _txt(data)
        if ct == "application/pdf" or fn.endswith(".pdf"):
            return _pdf(data)
        if fn.endswith(".docx") or "wordprocessingml" in ct:
            return _docx(data)
        if fn.endswith(".pptx") or "presentationml" in ct:
            return _pptx(data)
    except Exception as e:
        log.warning("attachment extract failed for %s (%s): %s", filename, content_type, e)
        return None
    return None  # unsupported type (image, zip, ...) — caller notes it
